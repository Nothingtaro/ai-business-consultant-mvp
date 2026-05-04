from __future__ import annotations

import logging
from dataclasses import dataclass
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from core.schemas import FinalConsultingReport
from utils.slide_content_adapter import DeckContent, SlideContent, SlideTable, adapt_report_for_slides


logger = logging.getLogger(__name__)

PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@dataclass(frozen=True)
class DeckDesign:
    slide_width = Inches(13.333)
    slide_height = Inches(7.5)

    margin_x = Inches(0.58)
    margin_top = Inches(0.55)
    content_top = Inches(1.22)
    footer_y = Inches(7.07)
    gutter = Inches(0.18)
    card_gap = Inches(0.22)

    font_body = "Aptos"
    font_display = "Aptos Display"
    font_title = 23
    font_subtitle = 8
    font_kicker = 7
    font_body_size = 9
    font_caption = 7
    font_table = 7
    font_table_header = 7
    min_body_font = 7

    ink = RGBColor(21, 32, 46)
    navy = RGBColor(20, 38, 61)
    slate = RGBColor(49, 63, 82)
    muted = RGBColor(96, 108, 123)
    light_muted = RGBColor(133, 145, 160)
    canvas = RGBColor(247, 249, 252)
    panel = RGBColor(255, 255, 255)
    border = RGBColor(213, 220, 229)
    table_alt = RGBColor(242, 246, 250)

    blue = RGBColor(38, 105, 213)
    teal = RGBColor(25, 145, 136)
    gold = RGBColor(195, 132, 42)
    red = RGBColor(183, 68, 75)
    violet = RGBColor(93, 91, 199)
    white = RGBColor(255, 255, 255)

    soft_blue = RGBColor(232, 241, 255)
    soft_teal = RGBColor(229, 245, 243)
    soft_gold = RGBColor(251, 244, 231)
    soft_red = RGBColor(252, 236, 237)
    soft_violet = RGBColor(239, 238, 252)

    footer_text = "AI Business Consulting Agent"


DESIGN = DeckDesign()
ACCENTS = [DESIGN.blue, DESIGN.teal, DESIGN.gold, DESIGN.violet, DESIGN.red]
SOFT_FILLS = [DESIGN.soft_blue, DESIGN.soft_teal, DESIGN.soft_gold, DESIGN.soft_violet, DESIGN.soft_red]


def build_pptx_report(report: FinalConsultingReport) -> bytes:
    content = adapt_report_for_slides(report)
    _log_qa_warnings(content)

    presentation = Presentation()
    presentation.slide_width = DESIGN.slide_width
    presentation.slide_height = DESIGN.slide_height

    add_title_slide(presentation, content.slides[0])
    add_section_divider_slide(presentation, content.slides[1])
    add_executive_summary_slide(presentation, content.slides[2])
    add_decision_context_slide(presentation, content.slides[3])
    add_issue_tree_slide(presentation, content.slides[4])
    add_hypotheses_slide(presentation, content.slides[5])
    add_options_comparison_slide(presentation, content.slides[6])
    add_recommendation_slide(presentation, content.slides[7])
    add_financials_slide(presentation, content.slides[8])
    add_risks_slide(presentation, content.slides[9])
    add_action_plan_slide(presentation, content.slides[10])
    add_deck_storyboard_slide(presentation, content.slides[11])
    add_critic_review_slide(presentation, content.slides[12])
    add_appendix_slide_if_needed(presentation, content)

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def add_title_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _blank_slide(presentation)
    _set_background(slide, DESIGN.canvas)

    _shape(slide, 0, 0, Inches(3.25), DESIGN.slide_height, DESIGN.navy, DESIGN.navy)
    _shape(slide, Inches(3.25), 0, Inches(0.09), DESIGN.slide_height, DESIGN.teal, DESIGN.teal)
    _shape(slide, Inches(3.34), Inches(0.0), Inches(9.99), Inches(0.18), DESIGN.blue, DESIGN.blue)

    _textbox(slide, "EXECUTIVE AI STRATEGY DECK", Inches(0.58), Inches(0.55), Inches(2.2), Inches(0.22), 7, DESIGN.white, bold=True)
    _textbox(slide, "Business\nDecision\nPackage", Inches(0.58), Inches(1.08), Inches(2.25), Inches(1.85), 27, DESIGN.white, bold=True)
    _textbox(slide, "Structured analysis for management review", Inches(0.58), Inches(6.62), Inches(2.15), Inches(0.42), 8, RGBColor(198, 209, 222))

    _textbox(slide, content.title, Inches(3.85), Inches(0.92), Inches(8.55), Inches(0.55), 28, DESIGN.ink, bold=True)
    _textbox(slide, content.payload["problem"], Inches(3.88), Inches(1.72), Inches(8.55), Inches(1.05), 16, DESIGN.muted)

    _kpi_card(slide, "Geography", content.payload["geography"], Inches(3.88), Inches(3.18), Inches(2.65), Inches(1.05), DESIGN.soft_blue)
    _kpi_card(slide, "Target Customers", content.payload["customers"], Inches(6.78), Inches(3.18), Inches(2.65), Inches(1.05), DESIGN.soft_teal)
    _kpi_card(slide, "Budget", content.payload["budget"], Inches(9.68), Inches(3.18), Inches(2.65), Inches(1.05), DESIGN.soft_gold)

    _callout(slide, "Expected Output", content.payload["expected_output"], Inches(3.88), Inches(4.82), Inches(8.45), Inches(1.08), DESIGN.teal)


def add_section_divider_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _blank_slide(presentation)
    _set_background(slide, DESIGN.navy)
    _shape(slide, 0, 0, Inches(0.16), DESIGN.slide_height, DESIGN.teal, DESIGN.teal)
    _shape(slide, Inches(0.16), Inches(0.0), Inches(0.06), DESIGN.slide_height, DESIGN.blue, DESIGN.blue)
    _textbox(slide, content.subtitle.upper(), Inches(0.85), Inches(0.82), Inches(5.5), Inches(0.28), 8, RGBColor(186, 214, 224), bold=True)
    _textbox(slide, content.title, Inches(0.85), Inches(1.55), Inches(6.8), Inches(0.75), 31, DESIGN.white, bold=True)
    _textbox(slide, content.key_message, Inches(0.88), Inches(2.55), Inches(7.2), Inches(0.85), 15, RGBColor(215, 224, 235))

    labels = content.payload.get("labels", [])
    left = Inches(0.88)
    for index, label in enumerate(labels):
        _process_step(slide, str(index + 1), label, left, Inches(4.78), Inches(2.65), Inches(0.85))
        left += Inches(2.95)
    _footer(slide, len(presentation.slides), dark=True)


def add_executive_summary_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    _hero_message(slide, "Key Message", content.key_message, DESIGN.margin_x, Inches(1.12), Inches(12.15), Inches(0.92), DESIGN.teal)

    cards = content.payload.get("cards", [])
    positions = [
        (DESIGN.margin_x, Inches(2.35)),
        (Inches(6.78), Inches(2.35)),
        (DESIGN.margin_x, Inches(4.58)),
        (Inches(6.78), Inches(4.58)),
    ]
    for index, ((title, body), (left, top)) in enumerate(zip(cards, positions)):
        _insight_card(slide, title, body, left, top, Inches(5.55), Inches(1.55), ACCENTS[index])


def add_decision_context_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    _hero_message(slide, "Decision Question", content.payload["decision"], DESIGN.margin_x, Inches(1.12), Inches(12.15), Inches(1.0), DESIGN.blue)
    _bullet_card(slide, "Business Context", content.payload["context"], DESIGN.margin_x, Inches(2.55), Inches(3.85), Inches(3.75), DESIGN.teal)
    _bullet_card(slide, "Constraints", content.payload["constraints"], Inches(4.72), Inches(2.55), Inches(2.55), Inches(3.75), DESIGN.gold)
    _bullet_card(slide, "Success Criteria", content.payload["success"], Inches(7.52), Inches(2.55), Inches(2.55), Inches(3.75), DESIGN.blue)
    _bullet_card(slide, "Key Unknowns", content.payload["unknowns"], Inches(10.32), Inches(2.55), Inches(2.48), Inches(3.75), DESIGN.red)


def add_issue_tree_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    branches = content.payload.get("branches", [])
    if not branches:
        _empty_state(slide)
        return

    positions = [
        (DESIGN.margin_x, Inches(1.22)),
        (Inches(6.78), Inches(1.22)),
        (DESIGN.margin_x, Inches(3.18)),
        (Inches(6.78), Inches(3.18)),
    ]
    for index, (branch, pos) in enumerate(zip(branches, positions)):
        _issue_branch_card(slide, branch, pos[0], pos[1], Inches(5.55), Inches(1.55), ACCENTS[index])
    _callout(slide, "Branch Logic", content.payload.get("branch_logic", content.key_message), DESIGN.margin_x, Inches(5.38), Inches(12.15), Inches(0.85), DESIGN.teal)


def add_hypotheses_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    table = content.payload.get("table")
    if not table or not table.headers:
        _empty_state(slide)
        return
    _table(slide, table, DESIGN.margin_x, Inches(1.18), Inches(12.15), Inches(4.35), col_widths=[Inches(4.45), Inches(3.7), Inches(4.0)])
    _callout(slide, "Initial Lean", content.payload.get("initial_lean", content.key_message), DESIGN.margin_x, Inches(5.85), Inches(12.15), Inches(0.68), DESIGN.gold)


def add_options_comparison_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    options = content.payload.get("options", [])
    if not options:
        _empty_state(slide)
        return

    left = DESIGN.margin_x
    for index, option in enumerate(options):
        _option_card(slide, index + 1, option, left, Inches(1.18), Inches(3.88), Inches(2.25), SOFT_FILLS[index], ACCENTS[index])
        left += Inches(4.14)

    rows = [(item["option"], item["upside"], item["downside"], item["implication"]) for item in options]
    _table(
        slide,
        SlideTable(["Option", "Upside", "Downside", "Implication"], rows),
        DESIGN.margin_x,
        Inches(3.82),
        Inches(12.15),
        Inches(2.62),
        col_widths=[Inches(2.4), Inches(3.0), Inches(3.0), Inches(3.75)],
        font_size=7,
    )


def add_recommendation_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    _hero_message(slide, "Recommended Action", content.payload["recommendation"], DESIGN.margin_x, Inches(1.12), Inches(12.15), Inches(1.02), DESIGN.teal)
    _bullet_card(slide, "Rationale", content.payload["rationale"], DESIGN.margin_x, Inches(2.55), Inches(5.85), Inches(3.75), DESIGN.blue)
    _callout(slide, "Financial Implication", content.payload["financial_implication"], Inches(6.75), Inches(2.55), Inches(6.05), Inches(1.35), DESIGN.gold)
    _bullet_card(slide, "Immediate Management Actions", content.payload["next_steps"], Inches(6.75), Inches(4.22), Inches(6.05), Inches(2.08), DESIGN.teal)


def add_financials_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    assumptions = content.payload.get("assumptions")
    scenarios = content.payload.get("scenarios")
    if not assumptions:
        _empty_state(slide)
        return

    _table(slide, assumptions, DESIGN.margin_x, Inches(1.18), Inches(12.15), Inches(2.45), col_widths=[Inches(4.35), Inches(2.35), Inches(2.35), Inches(3.1)], font_size=7)
    if scenarios:
        left = DESIGN.margin_x
        for index, row in enumerate(scenarios.rows[:3]):
            label, revenue, margin, profit = row
            _kpi_card(slide, label, profit, left, Inches(4.0), Inches(3.88), Inches(1.08), SOFT_FILLS[index], subtitle=f"Revenue {revenue} | GM {margin}")
            left += Inches(4.14)
    _callout(slide, "Break-Even / Scenario Logic", content.payload.get("break_even", content.key_message), DESIGN.margin_x, Inches(5.55), Inches(12.15), Inches(0.82), DESIGN.gold)


def add_risks_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    table = content.payload.get("table")
    if not table:
        _empty_state(slide)
        return
    _table(slide, table, DESIGN.margin_x, Inches(1.18), Inches(12.15), Inches(5.08), col_widths=[Inches(2.65), Inches(3.2), Inches(4.25), Inches(2.05)], font_size=7)
    _caption(slide, "Management note: confirm risk ownership during implementation planning.", DESIGN.margin_x, Inches(6.42), Inches(12.15))


def add_action_plan_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    columns = content.payload.get("columns", [])
    left = DESIGN.margin_x
    for index, (title, bullets) in enumerate(columns):
        _roadmap_column(slide, title, bullets, left, Inches(1.2), Inches(3.88), Inches(5.25), ACCENTS[index], SOFT_FILLS[index])
        left += Inches(4.14)


def add_deck_storyboard_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    table = content.payload.get("table")
    if not table:
        _empty_state(slide)
        return
    _table(slide, table, DESIGN.margin_x, Inches(1.08), Inches(12.15), Inches(5.72), col_widths=[Inches(0.5), Inches(2.55), Inches(6.8), Inches(2.3)], font_size=6.8)


def add_critic_review_slide(presentation: Presentation, content: SlideContent) -> None:
    slide = _content_slide(presentation, content)
    if not content.payload:
        _empty_state(slide)
        return

    _score_card(slide, content.payload["score"], DESIGN.margin_x, Inches(1.18), Inches(2.35), Inches(1.28))
    _hero_message(slide, "Final Verdict", content.key_message, Inches(3.15), Inches(1.18), Inches(9.65), Inches(1.28), DESIGN.blue)
    _bullet_card(slide, "Strengths", content.payload["strengths"], DESIGN.margin_x, Inches(2.85), Inches(3.85), Inches(3.5), DESIGN.teal)
    _bullet_card(slide, "Weaknesses", content.payload["weaknesses"], Inches(4.72), Inches(2.85), Inches(3.85), Inches(3.5), DESIGN.gold)
    _bullet_card(slide, "Critical Gaps", content.payload["gaps"], Inches(8.85), Inches(2.85), Inches(3.95), Inches(3.5), DESIGN.red)


def add_appendix_slide_if_needed(presentation: Presentation, content: DeckContent) -> None:
    for appendix_slide in content.appendix[:4]:
        slide = _content_slide(presentation, appendix_slide)
        if "table" in appendix_slide.payload:
            _table(slide, appendix_slide.payload["table"], DESIGN.margin_x, Inches(1.24), Inches(12.15), Inches(5.38), font_size=6.8)
        else:
            _bullet_card(slide, appendix_slide.key_message, appendix_slide.payload.get("bullets", []), DESIGN.margin_x, Inches(1.24), Inches(12.15), Inches(5.38), DESIGN.violet)


def _log_qa_warnings(content: DeckContent) -> None:
    for warning in content.warnings:
        logger.warning("PPTX export QA warning: %s", warning)


def _blank_slide(presentation: Presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _content_slide(presentation: Presentation, content: SlideContent):
    slide = _blank_slide(presentation)
    _set_background(slide, DESIGN.canvas)
    slide_number = len(presentation.slides)
    _shape(slide, 0, 0, DESIGN.slide_width, Inches(0.12), DESIGN.navy, DESIGN.navy)
    _shape(slide, 0, Inches(0.12), Inches(0.09), DESIGN.slide_height - Inches(0.12), DESIGN.teal, DESIGN.teal)
    _textbox(slide, "EXECUTIVE AI STRATEGY DECK", DESIGN.margin_x, Inches(0.28), Inches(2.8), Inches(0.18), DESIGN.font_kicker, DESIGN.blue, bold=True)
    _textbox(slide, content.title, DESIGN.margin_x, Inches(0.53), Inches(8.65), Inches(0.4), DESIGN.font_title, DESIGN.ink, bold=True)
    _textbox(slide, content.subtitle, Inches(9.35), Inches(0.58), Inches(3.1), Inches(0.22), DESIGN.font_subtitle, DESIGN.muted, align=PP_ALIGN.RIGHT)
    _footer(slide, slide_number)
    return slide


def _set_background(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None, *, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if rounded:
        shape.adjustments[0] = 0.08
    return shape


def _textbox(slide, text: object, left, top, width, height, font_size: float, color: RGBColor, *, bold: bool = False, align=PP_ALIGN.LEFT, max_chars: int = 420):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.text = _shorten(text, max_chars)
    paragraph.alignment = align
    paragraph.font.name = DESIGN.font_display if font_size >= 16 else DESIGN.font_body
    paragraph.font.size = Pt(max(font_size, DESIGN.min_body_font))
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _footer(slide, slide_number: int, *, dark: bool = False) -> None:
    color = RGBColor(181, 194, 208) if dark else DESIGN.light_muted
    _textbox(slide, DESIGN.footer_text, DESIGN.margin_x, DESIGN.footer_y, Inches(5.2), Inches(0.18), DESIGN.font_caption, color)
    _textbox(slide, f"{slide_number:02d}", Inches(12.32), DESIGN.footer_y, Inches(0.45), Inches(0.18), DESIGN.font_caption, color, align=PP_ALIGN.RIGHT)


def _caption(slide, text: str, left, top, width) -> None:
    _textbox(slide, text, left, top, width, Inches(0.18), DESIGN.font_caption, DESIGN.light_muted)


def _hero_message(slide, label: str, message: str, left, top, width, height, accent: RGBColor) -> None:
    _callout(slide, label, message, left, top, width, height, accent, body_size=13.5, max_chars=230)


def _callout(slide, title: str, body: object, left, top, width, height, accent: RGBColor, *, body_size: float = 10.2, max_chars: int = 170) -> None:
    shape = _shape(slide, left, top, width, height, DESIGN.panel, DESIGN.border, rounded=True)
    _shape(slide, left, top, Inches(0.11), height, accent, accent)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.3)
    frame.margin_right = Inches(0.18)
    frame.margin_top = Inches(0.1)
    frame.margin_bottom = Inches(0.08)
    frame.clear()
    title_para = frame.paragraphs[0]
    title_para.text = title.upper()
    _format_para(title_para, DESIGN.font_caption, accent, bold=True)
    body_para = frame.add_paragraph()
    body_para.text = _shorten(body, max_chars)
    _format_para(body_para, body_size, DESIGN.ink, bold=True, display=True)
    body_para.space_before = Pt(4)


def _insight_card(slide, title: str, body: str, left, top, width, height, accent: RGBColor) -> None:
    _shape(slide, left, top, width, height, DESIGN.panel, DESIGN.border, rounded=True)
    _shape(slide, left, top, Inches(0.08), height, accent, accent)
    _textbox(slide, title, left + Inches(0.24), top + Inches(0.16), width - Inches(0.42), Inches(0.24), 9, accent, bold=True, max_chars=70)
    _textbox(slide, body, left + Inches(0.24), top + Inches(0.55), width - Inches(0.42), height - Inches(0.65), 10.5, DESIGN.ink, bold=True, max_chars=155)


def _bullet_card(slide, title: str, bullets: list[str], left, top, width, height, accent: RGBColor) -> None:
    shape = _shape(slide, left, top, width, height, DESIGN.panel, DESIGN.border, rounded=True)
    _shape(slide, left, top, Inches(0.08), height, accent, accent)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.24)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.16)
    frame.margin_bottom = Inches(0.1)
    frame.clear()
    title_para = frame.paragraphs[0]
    title_para.text = title
    _format_para(title_para, 9.5, DESIGN.ink, bold=True)
    items = bullets[:5] or ["Not provided."]
    for item in items:
        para = frame.add_paragraph()
        para.text = f"- {_shorten(item, 118)}"
        _format_para(para, DESIGN.font_body_size, DESIGN.muted)
        para.space_before = Pt(5)


def _issue_branch_card(slide, branch: dict, left, top, width, height, accent: RGBColor) -> None:
    shape = _shape(slide, left, top, width, height, DESIGN.panel, DESIGN.border, rounded=True)
    _shape(slide, left, top, Inches(0.08), height, accent, accent)
    _textbox(slide, branch["name"], left + Inches(0.24), top + Inches(0.14), width - Inches(0.42), Inches(0.26), 10.5, DESIGN.ink, bold=True)
    _textbox(slide, branch["question"], left + Inches(0.24), top + Inches(0.5), width - Inches(0.42), Inches(0.38), 8.4, DESIGN.muted, max_chars=115)
    subs = " | ".join(branch.get("sub_branches", [])) or "Sub-branches TBD"
    _textbox(slide, subs, left + Inches(0.24), top + Inches(1.05), width - Inches(0.42), Inches(0.28), 7.6, accent, bold=True, max_chars=110)


def _option_card(slide, number: int, option: dict, left, top, width, height, fill: RGBColor, accent: RGBColor) -> None:
    _shape(slide, left, top, width, height, fill, DESIGN.border, rounded=True)
    _shape(slide, left + Inches(0.18), top + Inches(0.16), Inches(0.35), Inches(0.35), accent, accent, rounded=True)
    _textbox(slide, str(number), left + Inches(0.26), top + Inches(0.2), Inches(0.18), Inches(0.18), 8, DESIGN.white, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, option["option"], left + Inches(0.66), top + Inches(0.17), width - Inches(0.85), Inches(0.32), 11, DESIGN.ink, bold=True, max_chars=52)
    _textbox(slide, option["description"], left + Inches(0.24), top + Inches(0.68), width - Inches(0.48), Inches(0.58), 8.2, DESIGN.muted, max_chars=120)
    _textbox(slide, f"Implication: {option['implication']}", left + Inches(0.24), top + Inches(1.55), width - Inches(0.48), Inches(0.4), 8, accent, bold=True, max_chars=90)


def _roadmap_column(slide, title: str, bullets: list[str], left, top, width, height, accent: RGBColor, fill: RGBColor) -> None:
    _shape(slide, left, top, width, height, DESIGN.panel, DESIGN.border, rounded=True)
    _shape(slide, left, top, width, Inches(0.55), fill, DESIGN.border, rounded=True)
    _textbox(slide, title, left + Inches(0.2), top + Inches(0.17), width - Inches(0.4), Inches(0.22), 10.5, accent, bold=True)
    y = top + Inches(0.85)
    for index, item in enumerate((bullets or ["Not provided."])[:5]):
        _shape(slide, left + Inches(0.22), y + Inches(0.02), Inches(0.25), Inches(0.25), accent, accent, rounded=True)
        _textbox(slide, str(index + 1), left + Inches(0.29), y + Inches(0.055), Inches(0.1), Inches(0.1), 6.5, DESIGN.white, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, item, left + Inches(0.58), y, width - Inches(0.78), Inches(0.55), 8.6, DESIGN.muted, max_chars=105)
        y += Inches(0.78)


def _score_card(slide, score: str, left, top, width, height) -> None:
    fill = DESIGN.soft_teal if score.startswith(("4", "5")) else DESIGN.soft_gold if score.startswith("3") else DESIGN.soft_red
    _shape(slide, left, top, width, height, fill, DESIGN.border, rounded=True)
    _textbox(slide, score, left + Inches(0.2), top + Inches(0.18), width - Inches(0.4), Inches(0.44), 24, DESIGN.ink, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, "Decision Readiness", left + Inches(0.2), top + Inches(0.75), width - Inches(0.4), Inches(0.2), 8, DESIGN.muted, align=PP_ALIGN.CENTER)


def _kpi_card(slide, title: str, value: object, left, top, width, height, fill: RGBColor, *, subtitle: str = "") -> None:
    _shape(slide, left, top, width, height, fill, DESIGN.border, rounded=True)
    _textbox(slide, title.upper(), left + Inches(0.17), top + Inches(0.12), width - Inches(0.34), Inches(0.18), 7, DESIGN.muted, bold=True)
    _textbox(slide, value, left + Inches(0.17), top + Inches(0.38), width - Inches(0.34), Inches(0.36), 13, DESIGN.ink, bold=True, max_chars=80)
    if subtitle:
        _textbox(slide, subtitle, left + Inches(0.17), top + Inches(0.78), width - Inches(0.34), Inches(0.18), 6.8, DESIGN.muted, max_chars=80)


def _process_step(slide, number: str, label: str, left, top, width, height) -> None:
    _shape(slide, left, top, width, height, RGBColor(34, 55, 82), RGBColor(58, 84, 116), rounded=True)
    _textbox(slide, number, left + Inches(0.18), top + Inches(0.2), Inches(0.35), Inches(0.2), 10, DESIGN.teal, bold=True)
    _textbox(slide, label, left + Inches(0.62), top + Inches(0.2), width - Inches(0.85), Inches(0.22), 11, DESIGN.white, bold=True)


def _table(slide, table_content: SlideTable, left, top, width, height, *, col_widths: list | None = None, font_size: float | None = None) -> None:
    rows = table_content.rows or [tuple("Not provided." if index == 0 else "" for index in range(len(table_content.headers)))]
    headers = table_content.headers or ["Section"]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for index, col_width in enumerate(col_widths[: len(headers)]):
            table.columns[index].width = col_width

    header_h = Inches(0.34)
    table.rows[0].height = header_h
    row_h = int((height - header_h) / max(len(rows), 1))
    for index in range(1, len(rows) + 1):
        table.rows[index].height = row_h

    size = font_size or DESIGN.font_table
    for col_index, header in enumerate(headers):
        _cell(table.cell(0, col_index), header, DESIGN.navy, DESIGN.white, size, bold=True)
    for row_index, row in enumerate(rows, start=1):
        fill = DESIGN.panel if row_index % 2 else DESIGN.table_alt
        for col_index in range(len(headers)):
            value = row[col_index] if col_index < len(row) else ""
            _cell(table.cell(row_index, col_index), value, fill, DESIGN.muted, size)


def _cell(cell, text: object, fill: RGBColor, color: RGBColor, font_size: float, *, bold: bool = False) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.TOP
    frame = cell.text_frame
    frame.word_wrap = True
    frame.clear()
    para = frame.paragraphs[0]
    para.text = _shorten(text, 115)
    _format_para(para, max(font_size, DESIGN.min_body_font), color, bold=bold)


def _format_para(paragraph, font_size: float, color: RGBColor, *, bold: bool = False, display: bool = False) -> None:
    paragraph.font.name = DESIGN.font_display if display else DESIGN.font_body
    paragraph.font.size = Pt(max(font_size, DESIGN.min_body_font))
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _empty_state(slide) -> None:
    _callout(slide, "No Content Generated", "This section was not available in the final consulting report.", DESIGN.margin_x, Inches(1.3), Inches(12.15), Inches(0.9), DESIGN.gold)


def _shorten(value: object, max_chars: int = 150) -> str:
    text = str(value or "Not provided.").replace("\r", " ").replace("\n", " ").strip()
    text = " ".join(text.split())
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 3].rstrip() + "..."
